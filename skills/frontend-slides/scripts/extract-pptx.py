#!/usr/bin/env python3
"""
extract-pptx.py - Frontend Slides

Extract text, structure, images and speaker notes from a PowerPoint file
so the deck can be rebuilt as HTML.

    python scripts/extract-pptx.py <input.pptx> <output_dir>

Requires python-pptx:   pip install python-pptx

Writes:
    <output_dir>/content.json   full structured extraction
    <output_dir>/outline.md     human-readable outline for confirmation
    <output_dir>/assets/        every embedded image, named slideNN-imgN.ext
"""

import json
import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit(
        "error: python-pptx is not installed.\n"
        "  pip install python-pptx"
    )


def slugify(text, fallback="untitled"):
    s = re.sub(r"[^\w\s-]", "", (text or "").strip().lower())
    s = re.sub(r"[\s_-]+", "-", s).strip("-")
    return s or fallback


def shape_text(shape):
    """Return paragraphs as [{text, level, bold}], skipping empties."""
    if not shape.has_text_frame:
        return []
    out = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            continue
        bold = any(run.font.bold for run in para.runs if run.font.bold)
        out.append({"text": text, "level": para.level, "bold": bool(bold)})
    return out


def table_data(shape):
    rows = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def emu_to_px(value):
    """96 DPI - good enough for relative placement hints."""
    return round(Emu(value).inches * 96) if value is not None else None


def extract(pptx_path, out_dir):
    prs = Presentation(pptx_path)
    assets_dir = os.path.join(out_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    deck = {
        "source": os.path.basename(pptx_path),
        "slide_width_px": emu_to_px(prs.slide_width),
        "slide_height_px": emu_to_px(prs.slide_height),
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for idx, slide in enumerate(prs.slides, start=1):
        entry = {
            "index": idx,
            "layout": slide.slide_layout.name,
            "title": None,
            "blocks": [],
            "tables": [],
            "images": [],
            "notes": None,
        }

        # Title placeholder, when the layout defines one
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
            entry["title"] = title or None

        img_n = 0
        for shape in slide.shapes:
            # Skip the title shape - already captured
            if slide.shapes.title is not None and shape == slide.shapes.title:
                continue

            if shape.shape_type == 13 or getattr(shape, "image", None) is not None:
                try:
                    image = shape.image
                except (AttributeError, ValueError):
                    image = None
                if image is not None:
                    img_n += 1
                    ext = image.ext or "png"
                    name = "slide%02d-img%d.%s" % (idx, img_n, ext)
                    with open(os.path.join(assets_dir, name), "wb") as fh:
                        fh.write(image.blob)
                    entry["images"].append({
                        "file": "assets/" + name,
                        "width_px": emu_to_px(shape.width),
                        "height_px": emu_to_px(shape.height),
                        "left_px": emu_to_px(shape.left),
                        "top_px": emu_to_px(shape.top),
                        "alt": (shape.name or "").strip(),
                    })
                    continue

            if shape.has_table:
                entry["tables"].append(table_data(shape))
                continue

            paras = shape_text(shape)
            if paras:
                entry["blocks"].append({
                    "name": (shape.name or "").strip(),
                    "paragraphs": paras,
                })

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            entry["notes"] = notes or None

        # Fall back to the first line of text when there is no title placeholder
        if not entry["title"] and entry["blocks"]:
            entry["title"] = entry["blocks"][0]["paragraphs"][0]["text"][:80]

        deck["slides"].append(entry)

    # --- content.json ---
    with open(os.path.join(out_dir, "content.json"), "w", encoding="utf-8") as fh:
        json.dump(deck, fh, indent=2, ensure_ascii=False)

    # --- outline.md ---
    lines = [
        "# %s" % deck["source"],
        "",
        "%d slides - %sx%s px source canvas"
        % (deck["slide_count"], deck["slide_width_px"], deck["slide_height_px"]),
        "",
    ]
    for s in deck["slides"]:
        lines.append("## Slide %d - %s" % (s["index"], s["title"] or "(no title)"))
        lines.append("")
        lines.append("_Layout: %s_" % s["layout"])
        lines.append("")
        for block in s["blocks"]:
            for para in block["paragraphs"]:
                lines.append("%s- %s" % ("  " * para["level"], para["text"]))
        if s["tables"]:
            for t in s["tables"]:
                lines.append("")
                lines.append("Table (%d x %d): %s"
                             % (len(t), len(t[0]) if t else 0,
                                " | ".join(t[0]) if t else ""))
        if s["images"]:
            lines.append("")
            for im in s["images"]:
                lines.append("![%s](%s)" % (im["alt"] or "image", im["file"]))
        if s["notes"]:
            lines.append("")
            lines.append("> **Notes:** %s" % s["notes"].replace("\n", " "))
        lines.append("")

    with open(os.path.join(out_dir, "outline.md"), "w", encoding="utf-8") as fh:
        fh.write("\n".join(lines))

    total_images = sum(len(s["images"]) for s in deck["slides"])
    print("Extracted %d slides, %d images -> %s"
          % (deck["slide_count"], total_images, out_dir))
    print("  content.json  structured data")
    print("  outline.md    review this with the user before generating")
    print("  assets/       %d image file(s)" % total_images)


def main():
    if len(sys.argv) != 3:
        sys.exit("usage: python extract-pptx.py <input.pptx> <output_dir>")
    pptx_path, out_dir = sys.argv[1], sys.argv[2]
    if not os.path.isfile(pptx_path):
        sys.exit("error: no such file: %s" % pptx_path)
    os.makedirs(out_dir, exist_ok=True)
    extract(pptx_path, out_dir)


if __name__ == "__main__":
    main()
